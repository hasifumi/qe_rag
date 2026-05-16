"""テスト用架空就業規則文書を生成するスクリプト

生成ファイル:
  docs/就業規則説明.pptx  -- スライド6枚
  docs/福利厚生ガイド.pdf -- 福利厚生の詳細

実行:
  uv run python create_test_docs.py
"""
from pathlib import Path

DOCS_DIR = Path("./docs")
DOCS_DIR.mkdir(exist_ok=True)


# =========================================================
# 1. 就業規則説明.pptx
# =========================================================
def create_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slides_content = [
        {
            "title": "就業規則 2024年版",
            "body": "株式会社サンプル\n人事部",
        },
        {
            "title": "勤務時間",
            "body": (
                "■ 所定労働時間: 1日8時間（週40時間）\n"
                "■ 始業: 午前9時00分\n"
                "■ 終業: 午後6時00分\n"
                "■ 休憩: 正午12時〜午後1時（1時間）\n"
                "■ フレックスタイム制度あり（コアタイム 10:00〜15:00）"
            ),
        },
        {
            "title": "有給休暇",
            "body": (
                "■ 付与タイミング: 入社後6ヶ月経過後に10日付与\n"
                "■ 付与日数の増加: 勤続年数に応じて最大20日まで増加\n"
                "■ 時間単位取得: 1日を8分割して1時間単位で取得可能\n"
                "■ 繰越: 翌年度1年間に限り繰越可（最大40日）\n"
                "■ 取得手続き: 社内ポータルの「勤怠申請」メニューから申請"
            ),
        },
        {
            "title": "給与・賞与",
            "body": (
                "■ 給与形態: 月給制\n"
                "■ 支払日: 毎月25日（土日祝日の場合は前営業日）\n"
                "■ 残業代: 法定通り時間外割増賃金を別途支給\n"
                "■ 賞与: 年2回（6月・12月）、業績により支給\n"
                "■ 昇給: 年1回（4月）、人事考課に基づく"
            ),
        },
        {
            "title": "休日・休暇",
            "body": (
                "■ 週休2日制: 土曜日・日曜日\n"
                "■ 祝日: 国民の祝日に関する法律に定める日\n"
                "■ 年末年始: 12月29日〜1月3日\n"
                "■ 夏季休暇: 8月13日〜8月15日\n"
                "■ 慶弔休暇: 結婚5日、忌引き（続柄により1〜7日）\n"
                "■ 産前産後休業・育児休業: 法定通り"
            ),
        },
        {
            "title": "服務規律・その他",
            "body": (
                "■ テレワーク: 週2日まで在宅勤務可（要申請）\n"
                "■ 副業・兼業: 事前申請制（本業に支障がない範囲で許可）\n"
                "■ ハラスメント防止: 相談窓口は人事部またはコンプライアンス委員会\n"
                "■ 懲戒: 就業規則第8章に定める基準による"
            ),
        },
    ]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 完全に空白のレイアウト

    for i, content in enumerate(slides_content):
        slide = prs.slides.add_slide(blank_layout)

        # タイトルテキストボックス
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content["title"]
        p.font.size = Pt(28)
        p.font.bold = True

        # 本文テキストボックス
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content["body"]
        p.font.size = Pt(18)

    out = DOCS_DIR / "就業規則説明.pptx"
    prs.save(out)
    print(f"作成完了: {out}")


# =========================================================
# 2. 福利厚生ガイド.pdf
# =========================================================
def create_pdf():
    from fpdf import FPDF

    FONT_PATH = r"C:\Windows\Fonts\meiryo.ttc"
    if not Path(FONT_PATH).exists():
        # フォールバック: 日本語フォントが見つからない場合
        print(f"警告: フォントが見つかりません ({FONT_PATH})")
        print("  → Yu Gothic を試みます")
        FONT_PATH = r"C:\Windows\Fonts\YuGothR.ttc"
        if not Path(FONT_PATH).exists():
            raise FileNotFoundError(
                "日本語フォントが見つかりません。"
                "C:\\Windows\\Fonts\\ にある .ttc/.ttf フォントのパスを"
                "スクリプト内 FONT_PATH に手動で設定してください。"
            )

    sections = [
        {
            "heading": "福利厚生ガイド 2024年版",
            "body": "株式会社サンプル 人事部\n\n本ガイドでは社員向けの福利厚生制度を説明します。",
        },
        {
            "heading": "住宅手当・通勤手当",
            "body": (
                "【住宅手当】\n"
                "  月額20,000円を支給（賃貸契約書の提出が必要）。\n"
                "  持ち家の場合は対象外。社宅・寮入居者も対象外。\n\n"
                "【通勤手当】\n"
                "  公共交通機関の実費を支給（上限 月50,000円）。\n"
                "  マイカー通勤は距離に応じて支給（上限 月20,000円）。\n"
                "  申請は入社時および住所変更時に人事部へ提出。"
            ),
        },
        {
            "heading": "健康管理・医療支援",
            "body": (
                "【定期健康診断】\n"
                "  年1回実施（会社全額負担）。受診は義務。\n"
                "  35歳以上は生活習慣病予防健診（人間ドック相当）に変更可。\n\n"
                "【インフルエンザ予防接種】\n"
                "  毎年10〜11月に会社負担で実施（本人および家族1名まで）。\n\n"
                "【EAP（従業員支援プログラム）】\n"
                "  外部カウンセリングを無料で利用可（月2回まで）。"
            ),
        },
        {
            "heading": "育児・介護支援",
            "body": (
                "【育児休業】\n"
                "  子が2歳になるまで取得可能（法定の1歳6ヶ月を超えて延長可）。\n"
                "  パパ育休（産後パパ育休）制度あり: 子の出生後8週間以内に28日まで取得可。\n\n"
                "【育児短時間勤務】\n"
                "  子が小学3年生修了まで1日最大2時間の短縮勤務が可能。\n\n"
                "【介護休業】\n"
                "  対象家族1名につき通算93日取得可能。介護休暇（年5日）も別途あり。"
            ),
        },
        {
            "heading": "慶弔見舞金・その他給付",
            "body": (
                "【慶弔見舞金】\n"
                "  結婚祝い金: 50,000円\n"
                "  出産祝い金: 30,000円（第2子以降: 50,000円）\n"
                "  弔慰金（本人死亡）: 1,000,000円\n"
                "  弔慰金（配偶者・父母）: 100,000円\n\n"
                "【災害見舞金】\n"
                "  自然災害による損害に応じて50,000円〜300,000円を支給。\n\n"
                "【自己啓発支援】\n"
                "  業務関連の資格取得費用を年間100,000円まで補助。\n"
                "  社外研修・セミナー参加費用も別途申請可。"
            ),
        },
    ]

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_font("Meiryo", fname=FONT_PATH)

    for sec in sections:
        pdf.add_page()

        # 見出し
        pdf.set_font("Meiryo", size=16)
        pdf.set_fill_color(30, 80, 160)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 12, sec["heading"], fill=True, ln=True)
        pdf.ln(4)

        # 本文
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Meiryo", size=11)
        pdf.multi_cell(0, 7, sec["body"])

    out = DOCS_DIR / "福利厚生ガイド.pdf"
    pdf.output(str(out))
    print(f"作成完了: {out}")


if __name__ == "__main__":
    print("=== テスト文書生成 ===")
    print()

    print("① 就業規則説明.pptx を作成中...")
    create_pptx()

    print()
    print("② 福利厚生ガイド.pdf を作成中...")
    try:
        create_pdf()
    except ImportError:
        print("エラー: fpdf2 がインストールされていません。")
        print("  → pip install --user fpdf2  を実行してから再試行してください。")
    except FileNotFoundError as e:
        print(f"エラー: {e}")

    print()
    print("完了。docs/ フォルダを確認してください。")
