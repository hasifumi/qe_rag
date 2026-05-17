"""文書取り込み・ChromaDB登録モジュール"""
import os
import hashlib
from pathlib import Path
from datetime import datetime

import chromadb
from sentence_transformers import SentenceTransformer

import config


def _chunk_text(text: str, size: int = 1000, overlap: int = 200) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start : start + size])
        start += size - overlap
    return [c for c in chunks if c.strip()]


def _extract_pptx(path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(path)
    results = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join(
            shape.text for shape in slide.shapes if shape.has_text_frame
        )
        if text.strip():
            results.append({"text": text, "slide": i + 1})
    return results


def _extract_pdf(path: str) -> list[dict]:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    results = []
    for i, page in enumerate(extract_pages(path)):
        text = "".join(
            el.get_text() for el in page if isinstance(el, LTTextContainer)
        )
        if text.strip():
            results.append({"text": text, "page": i + 1})
    return results


def _extract_docx(path: str) -> list[dict]:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return [{"text": text, "page": 1}] if text.strip() else []


def _extract_text(path: str) -> list[dict]:
    text = Path(path).read_text(encoding="utf-8", errors="ignore")
    return [{"text": text, "page": 1}] if text.strip() else []


def _extract(path: str) -> list[dict]:
    ext = Path(path).suffix.lower()
    if ext == ".pptx":
        return _extract_pptx(path)
    elif ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext in (".txt", ".md"):
        return _extract_text(path)
    return []


def ingest(
    docs_dir: str = config.DOCS_DIR,
    db_dir: str = config.DB_DIR,
    on_progress=None,
) -> None:
    """文書を取り込み ChromaDB に登録する

    on_progress: 進捗メッセージを受け取る callable（引数1: str）。
                 None のとき print（CLI 後方互換）。Web UI は SSE 送出用に渡す。
    """
    report = on_progress if on_progress is not None else print
    embedder = SentenceTransformer(config.EMBED_MODEL)
    client = chromadb.PersistentClient(path=db_dir)
    col = client.get_or_create_collection("docs")

    doc_files = [
        p for p in Path(docs_dir).rglob("*")
        if p.suffix.lower() in (".pptx", ".pdf", ".docx", ".txt", ".md")
    ]
    report(f"文書ファイル数: {len(doc_files)}")

    # docs/ から削除されたファイルの ChromaDB エントリを除去
    present = {str(p) for p in doc_files}
    all_metas = col.get(include=["metadatas"])["metadatas"]
    registered = {m["file"] for m in all_metas if m.get("file")}
    for orphan in registered - present:
        col.delete(where={"file": orphan})
        report(f"  削除（docsから消失）: {Path(orphan).name}")

    for file_path in doc_files:
        path_str = str(file_path)
        mtime = str(os.path.getmtime(file_path))

        # 既登録かつ更新なしならスキップ
        existing = col.get(where={"file": path_str}, limit=1)
        if existing["ids"]:
            stored_mtime = existing["metadatas"][0].get("mtime", "")
            if stored_mtime == mtime:
                report(f"  スキップ（未変更）: {file_path.name}")
                continue

        # 既存エントリを削除して再登録
        col.delete(where={"file": path_str})

        sections = _extract(path_str)
        if not sections:
            report(f"  スキップ（テキストなし）: {file_path.name}")
            continue

        rel = file_path.relative_to(Path(docs_dir))
        department = rel.parts[0] if len(rel.parts) > 1 else "common"

        ids, docs, metas = [], [], []
        for sec in sections:
            for chunk in _chunk_text(sec["text"]):
                chunk_id = hashlib.md5(
                    (path_str + chunk).encode()
                ).hexdigest()
                ids.append(chunk_id)
                docs.append(chunk)
                meta = {"file": path_str, "mtime": mtime, "department": department}
                if "slide" in sec:
                    meta["slide"] = sec["slide"]
                elif "page" in sec:
                    meta["page"] = sec["page"]
                metas.append(meta)

        if ids:
            embeddings = embedder.encode(docs).tolist()
            col.add(ids=ids, documents=docs, embeddings=embeddings, metadatas=metas)
            report(f"  登録: {file_path.name} ({len(ids)} チャンク)")

    report("インデックス構築完了")
